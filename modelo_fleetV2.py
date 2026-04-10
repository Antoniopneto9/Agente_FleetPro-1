import os
import glob
import re
import hashlib
import shutil
from typing import Optional

import streamlit as st

st.set_page_config(page_title="FleetPro Expert", layout="wide", initial_sidebar_state="expanded")

try:
    import extra_streamlit_components as stx
    _cookie_manager = stx.CookieManager(key="fleetpro_cookies")
    _COOKIES_OK = True
except Exception:
    _cookie_manager = None
    _COOKIES_OK = False


def safe_run(fn):
    try:
        fn()
    except Exception as e:
        st.error("O app falhou ao inicializar. O erro também deve aparecer no terminal.")
        st.exception(e)
        st.stop()


# Dependências obrigatórias
try:
    import pandas as pd
except Exception as e:
    st.error("Dependência ausente: pandas")
    st.exception(e)
    st.stop()

try:
    import openpyxl  # noqa: F401
except Exception as e:
    st.error("Dependência ausente: openpyxl (necessário para ler .xlsx)")
    st.exception(e)
    st.stop()

# ConversationBufferMemory definido acima como classe própria

# Implementação própria de memória — sem depender do langchain.memory
class ConversationBufferMemory:
    """Memória de conversa simples, sem dependência do langchain.memory."""
    def __init__(self):
        self._msgs = []
        self.chat_memory = self

    def add_user_message(self, m):
        self._msgs.append(("human", m))

    def add_ai_message(self, m):
        self._msgs.append(("ai", m))

    @property
    def buffer_as_messages(self):
        return [type("Msg", (), {"type": t, "content": c})() for t, c in self._msgs]

from langchain_groq import ChatGroq

try:
    from langchain_openai import ChatOpenAI, OpenAIEmbeddings
except ImportError:
    ChatOpenAI = None
    OpenAIEmbeddings = None

try:
    from langchain_text_splitters import RecursiveCharacterTextSplitter
except ImportError:
    try:
        from langchain.text_splitter import RecursiveCharacterTextSplitter
    except ImportError:
        from langchain_community.text_splitter import RecursiveCharacterTextSplitter

from langchain_community.vectorstores import Chroma
from langchain_community.document_loaders import PyPDFLoader, TextLoader, CSVLoader
from langchain_core.documents import Document

# ======================
# Config / Paths
# ======================
APP_DIR = os.path.dirname(os.path.abspath(__file__))

BASE_DOCS_DIR = os.path.join(APP_DIR, "base_docs")

# Streamlit Cloud tem filesystem read-only — usa /tmp para o índice vetorial
# Localmente usa a pasta do projeto para persistir entre sessões
# Detecta ambiente cloud — usa /tmp que é sempre gravável
_IS_CLOUD = os.path.exists("/mount/src") or not os.access(APP_DIR, os.W_OK)
if _IS_CLOUD:
    CHROMA_DIR = "/tmp/fleetpro_chroma_db"
    os.makedirs(CHROMA_DIR, exist_ok=True)
else:
    CHROMA_DIR = os.path.join(APP_DIR, "chroma_db")
COLLECTION_NAME = "fleetpro_kb_v2"

# Arquivo principal de lookup de PN (CSV preferido, Excel como fallback)
MATRIX_CSV   = "Matriz_FP.csv"
MATRIX_EXCEL = "Matriz_FP.xlsx"
SHEET_FP_MATRIZ = "FP MATRIZ"
COLUNAS_BUSCA_PN = ["pn_gen", "pn_alternative", "pn_nxp", "pn_fleetpro"]

# Colunas que contêm modelos de equipamentos CASE IH / NHAG (texto com nomes de modelos)
COLUNAS_MODELOS_EQUIP = [
    "tractor_case_ih", "combine_case_ih", "headers_case_ih",
    "sch_case_ih", "sprayers_case_ih", "planters_case_ih",
    "other_machines_case_ih",
    "tractor_nhag", "combine_nhag", "headers_nhag",
    "sprayers_nhag", "planters_nhag",
    "forage_balers_and_others_nhag", "other_machines_nhag",
]

# Colunas que contêm PNs originais de outras marcas
COLUNAS_MARCAS_PN = ["john_deere", "macdon", "agco", "ideal", "massey_ferguson", "valtra"]

# Coluna de projeto de marketing (busca por tipo de produto/categoria)
COLUNA_MARKETING = "marketing_project"

# Valores exatos que aparecem na coluna MARKETING PROJECT da planilha.
# Cada entrada é um valor real + sinônimos/variações que o usuário pode digitar.
# A busca é feita com str.contains() parcial, então "ROLAMENTO" bate em "ROLAMENTOS" etc.
PALAVRAS_MARKETING = [
    # ── Valores exatos da coluna ─────────────────────────────────────────────
    "BUCKET",
    "ALL MAKES",
    "UNDERCARRIAGE",
    "PERIODICAL MAINTENANCE KITS",
    "CORREIAS",
    "FILTROS",
    "ROLAMENTOS",
    "PEÇAS DE DESGASTES",
    "PINOS & BUCHAS",
    "PLÁSTICOS",
    "CILINDROS",
    "LUBRIFICANTES",
    "CORRENTES",
    "VEDAÇÃO",
    "PULVERIZAÇÃO",
    "USINADOS",
    # ── Sinônimos / variações em português ───────────────────────────────────
    "CORREIA",          # → CORREIAS
    "FILTRO",           # → FILTROS
    "ROLAMENTO",        # → ROLAMENTOS
    "PECA DE DESGASTE", # → PEÇAS DE DESGASTES
    "PEÇA DE DESGASTE",
    "DESGASTE",
    "PINO",             # → PINOS & BUCHAS
    "PINOS",
    "BUCHA",            # → PINOS & BUCHAS
    "BUCHAS",
    "PLASTICO",         # → PLÁSTICOS
    "PLASTICOS",
    "PLÁSTICO",
    "CILINDRO",         # → CILINDROS
    "LUBRIFICANTE",     # → LUBRIFICANTES
    "CORRENTE",         # → CORRENTES
    "VEDACAO",          # → VEDAÇÃO (sem acento)
    "PULVERIZACAO",     # → PULVERIZAÇÃO (sem acento)
    "USINADO",          # → USINADOS
    "KIT",              # → PERIODICAL MAINTENANCE KITS
    "KITS",
    "MANUTENCAO",       # → PERIODICAL MAINTENANCE KITS
    "MANUTENÇÃO",
    # ── Sinônimos em inglês ───────────────────────────────────────────────────
    "BELT",             # → CORREIAS
    "BELTS",
    "FILTER",           # → FILTROS
    "FILTERS",
    "BEARING",          # → ROLAMENTOS
    "BEARINGS",
    "BUSHING",          # → PINOS & BUCHAS
    "BUSHINGS",
    "PIN",              # → PINOS & BUCHAS
    "PINS",
    "CYLINDER",         # → CILINDROS
    "CYLINDERS",
    "LUBRICANT",        # → LUBRIFICANTES
    "LUBRICANTS",
    "CHAIN",            # → CORRENTES
    "CHAINS",
    "SEAL",             # → VEDAÇÃO
    "SEALS",
    "SPRAY",            # → PULVERIZAÇÃO
    "SPRAYING",
    "WEAR",             # → PEÇAS DE DESGASTES
    "WEAR PARTS",
    "MAINTENANCE",      # → PERIODICAL MAINTENANCE KITS
]

# Mapeia sinônimos/variações → valor EXATO que aparece na coluna MARKETING PROJECT.
# Se o termo detectado não está aqui, ele é usado diretamente como filtro parcial.
MAPA_SINONIMOS_MARKETING = {
    # Português singular → plural (valor da coluna)
    "CORREIA":          "CORREIAS",
    "FILTRO":           "FILTROS",
    "ROLAMENTO":        "ROLAMENTOS",
    "PECA DE DESGASTE": "PEÇAS DE DESGASTES",
    "PEÇA DE DESGASTE": "PEÇAS DE DESGASTES",
    "DESGASTE":         "PEÇAS DE DESGASTES",
    "PINO":             "PINOS & BUCHAS",
    "PINOS":            "PINOS & BUCHAS",
    "BUCHA":            "PINOS & BUCHAS",
    "BUCHAS":           "PINOS & BUCHAS",
    "PLASTICO":         "PLÁSTICOS",
    "PLASTICOS":        "PLÁSTICOS",
    "PLÁSTICO":         "PLÁSTICOS",
    "CILINDRO":         "CILINDROS",
    "LUBRIFICANTE":     "LUBRIFICANTES",
    "CORRENTE":         "CORRENTES",
    "VEDACAO":          "VEDAÇÃO",
    "PULVERIZACAO":     "PULVERIZAÇÃO",
    "USINADO":          "USINADOS",
    "KIT":              "PERIODICAL MAINTENANCE KITS",
    "KITS":             "PERIODICAL MAINTENANCE KITS",
    "MANUTENCAO":       "PERIODICAL MAINTENANCE KITS",
    "MANUTENÇÃO":       "PERIODICAL MAINTENANCE KITS",
    "MAINTENANCE":      "PERIODICAL MAINTENANCE KITS",
    # Inglês → valor da coluna
    "BELT":             "CORREIAS",
    "BELTS":            "CORREIAS",
    "FILTER":           "FILTROS",
    "FILTERS":          "FILTROS",
    "BEARING":          "ROLAMENTOS",
    "BEARINGS":         "ROLAMENTOS",
    "BUSHING":          "PINOS & BUCHAS",
    "BUSHINGS":         "PINOS & BUCHAS",
    "PIN":              "PINOS & BUCHAS",
    "PINS":             "PINOS & BUCHAS",
    "CYLINDER":         "CILINDROS",
    "CYLINDERS":        "CILINDROS",
    "LUBRICANT":        "LUBRIFICANTES",
    "LUBRICANTS":       "LUBRIFICANTES",
    "CHAIN":            "CORRENTES",
    "CHAINS":           "CORRENTES",
    "SEAL":             "VEDAÇÃO",
    "SEALS":            "VEDAÇÃO",
    "SPRAY":            "PULVERIZAÇÃO",
    "SPRAYING":         "PULVERIZAÇÃO",
    "WEAR":             "PEÇAS DE DESGASTES",
    "WEAR PARTS":       "PEÇAS DE DESGASTES",
}


# Mapa de palavras-chave para nomes de colunas (busca por fabricante/tipo)
MAPA_PALAVRAS_COLUNAS = {
    "JOHN DEERE": ["john_deere"],
    "JD": ["john_deere"],
    "DEERE": ["john_deere"],
    "MACDON": ["macdon"],
    "AGCO": ["agco"],
    "IDEAL": ["ideal"],
    "MASSEY": ["massey_ferguson"],
    "MASSEY FERGUSON": ["massey_ferguson"],
    "MF": ["massey_ferguson"],
    "VALTRA": ["valtra"],
    "CASE": ["tractor_case_ih", "combine_case_ih", "headers_case_ih",
             "sch_case_ih", "sprayers_case_ih", "planters_case_ih",
             "other_machines_case_ih"],
    "CASE IH": ["tractor_case_ih", "combine_case_ih", "headers_case_ih",
                "sch_case_ih", "sprayers_case_ih", "planters_case_ih",
                "other_machines_case_ih"],
    "TRATOR": ["tractor_case_ih", "tractor_nhag"],
    "TRACTOR": ["tractor_case_ih", "tractor_nhag"],
    "COLHEITADEIRA": ["combine_case_ih", "combine_nhag"],
    "COMBINE": ["combine_case_ih", "combine_nhag"],
    "HEADER": ["headers_case_ih", "headers_nhag"],
    "PLATAFORMA": ["headers_case_ih", "headers_nhag"],
    "PULVERIZADOR": ["sprayers_case_ih", "sprayers_nhag"],
    "SPRAYER": ["sprayers_case_ih", "sprayers_nhag"],
    "PLANTADEIRA": ["planters_case_ih", "planters_nhag"],
    "PLANTER": ["planters_case_ih", "planters_nhag"],
    "NHAG": ["tractor_nhag", "combine_nhag", "headers_nhag",
             "sprayers_nhag", "planters_nhag",
             "forage_balers_and_others_nhag", "other_machines_nhag"],
    "NEW HOLLAND": ["tractor_nhag", "combine_nhag", "headers_nhag",
                    "sprayers_nhag", "planters_nhag",
                    "forage_balers_and_others_nhag", "other_machines_nhag"],
    # ── Modelos específicos Case IH ──────────────────────────────────────
    "MAGNUM":     ["tractor_case_ih"],
    "MAXXUM":     ["tractor_case_ih"],
    "OPTUM":      ["tractor_case_ih"],
    "FARMALL":    ["tractor_case_ih"],
    "PUMA":       ["tractor_case_ih"],
    "AXIAL":      ["combine_case_ih"],
    "AXIAL-FLOW": ["combine_case_ih"],
    "AF":         ["combine_case_ih"],
    "DRAPER":     ["headers_case_ih"],
    "FLEX":       ["headers_case_ih"],
    "PATRIOT":    ["sprayers_case_ih"],
    "EARLY RISER":["planters_case_ih"],
    # ── Modelos específicos New Holland ──────────────────────────────────
    "T6":         ["tractor_nhag"],
    "T7":         ["tractor_nhag"],
    "T8":         ["tractor_nhag"],
    "T9":         ["tractor_nhag"],
    "TC":         ["combine_nhag"],
    "CR":         ["combine_nhag"],
    "CX":         ["combine_nhag"],
    "FR":         ["forage_balers_and_others_nhag"],
    "GUARDIAN":   ["sprayers_nhag"],
    "BOOMER":     ["tractor_nhag"],
    # Segmentos
    "AG": ["ag"],
    "AGRICULTURE": ["ag"],
    "AGRICOLA": ["ag"],
    "AGRÍCOLA": ["ag"],
    "CE": ["ce"],
    "CONSTRUCTION": ["ce"],
    "CONSTRUCAO": ["ce"],
    "CONSTRUÇÃO": ["ce"],
    "AM": ["am"],
}

TIPOS_RAG = (".pdf", ".txt", ".csv", ".pptx", ".ppt", ".md", ".docx")

CONFIG_MODELOS = {
    "Groq": {
        "modelos": [
            "llama-3.3-70b-versatile",
            "llama-3.1-8b-instant",
        ],
        "chat": ChatGroq,
    },
    "OpenAI": {
        "modelos": ["gpt-4o-mini", "gpt-4o"],
        "chat": ChatOpenAI,
    },
}


# ======================
# Utilitários
# ======================
def norm_pn(x) -> str:
    """Normaliza o Part Number: maiúsculo, sem espaços, hífens e pontos.
    Trata floats do Excel: '47390059.0' → '47390059' (remove .0 antes de normalizar).
    """
    if x is None:
        return ""
    s = str(x).strip()
    # Remove sufixo .0 gerado pelo Excel ao ler números como float
    if re.match(r"^\d+\.0$", s):
        s = s[:-2]
    s = s.upper()
    s = re.sub(r"\s+", "", s)
    s = s.replace("-", "").replace(".", "")
    return s


_STOPWORDS_PN = {
    "ITEM", "PECA", "PEÇA", "NUMERO", "NÚMERO", "GENUINO", "GENUÍNO",
    "ORIGINAL", "SUBSTITUTO", "CODIGO", "CÓDIGO", "AJUDE", "VEJA",
    "OPCOES", "OPÇÕES", "ESSE", "ESTE", "PARA", "QUAL", "QUAIS",
    "BUSCA", "ENCONTRA", "ACHAR", "TENHO", "PRECISO", "PRODUCT",
    "NUMBER", "PART", "FIND", "SEARCH", "HELP", "THE", "AND", "PN",
    "COM", "UMA", "VER", "NAO", "NÃO", "SIM", "POR", "QUE", "MEU",
    "MIM", "ELE", "ELA", "SER", "TEM", "FOI", "TIPO", "MODELO",
}


def extrair_pns_da_mensagem(texto: str) -> list:
    candidatos_raw = re.findall(r"[A-Za-z0-9]+(?:[-\.][A-Za-z0-9]+)*", texto)
    candidatos = []
    for c in candidatos_raw:
        norm = norm_pn(c)
        if (
            len(norm) >= 4
            and re.search(r"\d", norm)
            and norm not in _STOPWORDS_PN
            and not norm.isalpha()
        ):
            candidatos.append(norm)
    return sorted(set(candidatos), key=lambda x: -len(x))


def _split_docs(docs):
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=800,         # reduzido para caber no limite de tokens do Groq
        chunk_overlap=100,
        separators=[
            "\n### Slide",    # respeita divisão entre slides do .md
            "\n\n",
            "\n",
            ".", ";", ",", " ", "",
        ],
    )
    return splitter.split_documents(docs)


# ======================
# RAG – Documentos de conhecimento
# ======================
# Arquivos a excluir do RAG (são usados pelo lookup direto ou são binários)
_EXCLUIR_DO_RAG = {"Matriz_FP.csv", "Matriz_FP.xlsx", "Matriz_FP.xls"}

def _listar_arquivos_rag(pasta: str):
    """
    Lista recursivamente arquivos RAG em subpastas de base_docs/.
    EXCLUI: arquivos da raiz (logos, Excel/CSV da matriz) — apenas subpastas entram no RAG.
    """
    arquivos = []
    for ext in TIPOS_RAG:
        for arq in glob.glob(os.path.join(pasta, "**", f"*{ext}"), recursive=True):
            if not os.path.isfile(arq):
                continue
            nome = os.path.basename(arq)
            # Exclui arquivos da matriz e arquivos direto na raiz base_docs
            if nome in _EXCLUIR_DO_RAG:
                continue
            # Exclui arquivos direto na raiz (logos, etc) — só subpastas entram
            if os.path.dirname(arq) == pasta:
                continue
            arquivos.append(arq)
    return sorted(set(arquivos))


def _hash_dos_arquivos(lista_arquivos) -> str:
    h = hashlib.sha256()
    for path in lista_arquivos:
        stat = os.stat(path)
        h.update(path.encode("utf-8"))
        h.update(str(stat.st_mtime).encode("utf-8"))
        h.update(str(stat.st_size).encode("utf-8"))
    return h.hexdigest()


def _hash_completo(lista_arquivos: list) -> str:
    """Hash combinando arquivos locais + flag do site."""
    h = hashlib.sha256()
    for path in lista_arquivos:
        stat = os.stat(path)
        h.update(path.encode("utf-8"))
        h.update(str(stat.st_mtime).encode("utf-8"))
        h.update(str(stat.st_size).encode("utf-8"))
    return h.hexdigest()


def _ler_hash_salvo():
    hash_path = os.path.join(CHROMA_DIR, "base_hash.txt")
    if os.path.exists(hash_path):
        with open(hash_path, "r", encoding="utf-8") as f:
            return f.read().strip()
    return None


def _salvar_hash(base_hash: str):
    os.makedirs(CHROMA_DIR, exist_ok=True)
    hash_path = os.path.join(CHROMA_DIR, "base_hash.txt")
    with open(hash_path, "w", encoding="utf-8") as f:
        f.write(base_hash)


def _apagar_indice():
    if os.path.isdir(CHROMA_DIR):
        shutil.rmtree(CHROMA_DIR, ignore_errors=True)


def _carregar_pptx(path: str) -> list:
    """Extrai texto de arquivos .pptx/.ppt slide a slide."""
    try:
        from pptx import Presentation
    except ImportError:
        st.warning("python-pptx não instalado. Instale com: pip install python-pptx")
        return []
    try:
        prs = Presentation(path)
        docs_ppt = []
        nome = os.path.basename(path)
        pasta_rel = os.path.relpath(os.path.dirname(path))
        for i, slide in enumerate(prs.slides, start=1):
            textos = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    textos.append(shape.text.strip())
            if textos:
                conteudo = "[Arquivo: " + str(nome) + " | Pasta: " + str(pasta_rel) + " | Slide " + str(i) + "\n" + "\n".join(textos)
                docs_ppt.append(Document(
                    page_content=conteudo,
                    metadata={"source": path, "slide": i, "tipo": "pptx"},
                ))
        return docs_ppt
    except Exception as e:
        st.warning(f"Não consegui ler {os.path.basename(path)}: {e}")
        return []


def _deve_indexar(path: str, pasta_raiz: str) -> bool:
    """Retorna True se o arquivo deve ser indexado no RAG."""
    nome = os.path.basename(path)
    # Exclui arquivos da matriz
    if nome in _EXCLUIR_DO_RAG:
        return False
    # Exclui arquivos direto na raiz base_docs (logos, etc)
    if os.path.dirname(os.path.abspath(path)) == os.path.abspath(pasta_raiz):
        return False
    return True


def _carregar_documentos_rag(pasta: str):
    """
    Carrega documentos RAG apenas das SUBPASTAS de base_docs/.
    Exclui: Matriz_FP.csv/xlsx, arquivos na raiz, imagens.
    """
    docs = []

    for path in glob.glob(os.path.join(pasta, "**", "*.pdf"), recursive=True):
        if not _deve_indexar(path, pasta):
            continue
        try:
            docs.extend(PyPDFLoader(path).load())
        except Exception as e:
            st.warning(f"Não consegui ler {os.path.basename(path)}: {e}")

    for path in glob.glob(os.path.join(pasta, "**", "*.txt"), recursive=True):
        if not _deve_indexar(path, pasta):
            continue
        try:
            docs.extend(TextLoader(path, encoding="utf-8").load())
        except Exception as e:
            st.warning(f"Não consegui ler {os.path.basename(path)}: {e}")

    for path in glob.glob(os.path.join(pasta, "**", "*.csv"), recursive=True):
        if not _deve_indexar(path, pasta):
            continue
        try:
            docs.extend(CSVLoader(path, encoding="utf-8").load())
        except Exception as e:
            st.warning(f"Não consegui ler {os.path.basename(path)}: {e}")

    for path in glob.glob(os.path.join(pasta, "**", "*.pptx"), recursive=True):
        if not _deve_indexar(path, pasta):
            continue
        docs.extend(_carregar_pptx(path))

    for path in glob.glob(os.path.join(pasta, "**", "*.ppt"), recursive=True):
        if not _deve_indexar(path, pasta):
            continue
        docs.extend(_carregar_pptx(path))

    for path in glob.glob(os.path.join(pasta, "**", "*.md"), recursive=True):
        if not _deve_indexar(path, pasta):
            continue
        try:
            md_docs = TextLoader(path, encoding="utf-8").load()
            # Injeta nome do arquivo e pasta no conteúdo para melhorar busca semântica
            nome_arquivo = os.path.splitext(os.path.basename(path))[0].replace("_", " ")
            pasta_nome = os.path.basename(os.path.dirname(path))
            for doc in md_docs:
                doc.page_content = f"[Fonte: {nome_arquivo} | Pasta: {pasta_nome}]\n" + doc.page_content
                doc.metadata["nome_arquivo"] = nome_arquivo
                doc.metadata["pasta"] = pasta_nome
            docs.extend(md_docs)
        except Exception as e:
            st.warning(f"Não consegui ler {os.path.basename(path)}: {e}")

    for path in glob.glob(os.path.join(pasta, "**", "*.docx"), recursive=True):
        if not _deve_indexar(path, pasta):
            continue
        try:
            from langchain_community.document_loaders import Docx2txtLoader
            docs.extend(Docx2txtLoader(path).load())
        except Exception as e:
            st.warning(f"Não consegui ler {os.path.basename(path)}: {e}")

    return docs


def _construir_vectorstore():
    """Cria um vectorstore novo do zero. Retorna None se embeddings não disponíveis."""
    os.makedirs(BASE_DOCS_DIR, exist_ok=True)
    try:
        from langchain_huggingface import HuggingFaceEmbeddings
    except ImportError:
        from langchain_community.embeddings import HuggingFaceEmbeddings
    try:
        embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
        # Testa se o modelo está disponível localmente
        embeddings.embed_query("teste")
    except Exception:
        st.session_state["_rag_indisponivel"] = True
        return None

    st.session_state["_rag_indisponivel"] = False
    docs = _carregar_documentos_rag(BASE_DOCS_DIR)
    if not docs:
        return None
    chunks = _split_docs(docs)
    if _IS_CLOUD:
        os.makedirs(CHROMA_DIR, exist_ok=True)
        return Chroma.from_documents(documents=chunks, embedding=embeddings, collection_name=COLLECTION_NAME, persist_directory=CHROMA_DIR)
    if os.path.isdir(CHROMA_DIR):
        shutil.rmtree(CHROMA_DIR, ignore_errors=True)
    os.makedirs(CHROMA_DIR, exist_ok=True)
    vs = Chroma.from_documents(documents=chunks, embedding=embeddings, collection_name=COLLECTION_NAME, persist_directory=CHROMA_DIR)
    _salvar_hash(_hash_completo(_listar_arquivos_rag(BASE_DOCS_DIR)))
    return vs


def obter_vectorstore():
    """Retorna o vectorstore do session_state, construindo se necessário."""
    if st.session_state.get("_vectorstore") is None:
        vs = _construir_vectorstore()
        st.session_state["_vectorstore"] = vs
    return st.session_state["_vectorstore"]


def resetar_vectorstore():
    """Descarta o vectorstore em cache para forçar reconstrução na próxima chamada."""
    st.session_state["_vectorstore"] = None
    st.session_state["_rag_indisponivel"] = False

# Mapa de siglas/termos → fragmentos de nome de arquivo para filtro
_MAPA_SIGLAS_ARQUIVO = {
    "VV":        ["VV"],
    "TEEJET":    ["TEEJET"],
    "REXNORD":   ["REX", "CORRENTES"],
    "REX":       ["REX", "CORRENTES"],
    "CORRENTE":  ["REX", "CORRENTES"],
    "CORRENTES": ["REX", "CORRENTES"],
    "IPESA":     ["IPESA"],
    "PETRONAS":  ["PETRONAS"],
    "FLEETPRO DAY": ["FLEETPRO DAY"],
    "DAY":       ["FLEETPRO DAY"],
}


def _detectar_filtro_arquivo(query: str) -> list:
    """Retorna lista de fragmentos para filtrar por nome de arquivo, ou [] se nenhum detectado."""
    q = query.upper()
    for sigla, fragmentos in _MAPA_SIGLAS_ARQUIVO.items():
        if re.search(r"\b" + re.escape(sigla) + r"\b", q):
            return fragmentos
    return []


def buscar_no_rag(vectorstore, query: str, k: int = 3):
    """
    Retorna (contexto_str, lista_de_fontes).
    Detecta siglas na query e filtra por nome de arquivo quando possível.
    k=3 para caber no limite de tokens do Groq free (12k TPM).
    """
    if vectorstore is None:
        return "", []

    try:
        # Busca semântica padrão
        resultados = vectorstore.similarity_search(query, k=k)

        # ── Filtro por sigla: se query menciona VV, TEEJET etc,
        #    busca mais resultados e filtra pelos que batem no nome do arquivo ──
        filtros = _detectar_filtro_arquivo(query)
        if filtros:
            candidatos = vectorstore.similarity_search(query, k=20)
            filtrados = [
                doc for doc in candidatos
                if any(
                    f.upper() in doc.metadata.get("source", "").upper()
                    or f.upper() in doc.metadata.get("nome_arquivo", "").upper()
                    for f in filtros
                )
            ]
            if filtrados:
                # Usa os filtrados (até k) + completa com semântica se necessário
                resultados = filtrados[:k]
                if len(resultados) < k:
                    for doc in candidatos:
                        if doc not in resultados:
                            resultados.append(doc)
                        if len(resultados) >= k:
                            break

        if not resultados:
            return "", []

        trechos = []
        fontes = []
        labels_vistos = set()

        for i, doc in enumerate(resultados, 1):
            fonte = doc.metadata.get("source", "documento")
            tipo = doc.metadata.get("tipo", "")
            slide = doc.metadata.get("slide", None)

            nome = doc.metadata.get("nome_arquivo", os.path.splitext(os.path.basename(fonte))[0])
            pasta = doc.metadata.get("pasta", "")
            label = "📄 " + nome + (" [" + pasta + "]" if pasta else "")
            fonte_url = fonte

            trecho_texto = doc.page_content.strip()[:600]
            trechos.append("[Trecho " + str(i) + " – " + label + "]\n" + trecho_texto)
            if label not in labels_vistos:
                labels_vistos.add(label)
                fontes.append({"label": label, "path": fonte_url, "tipo": tipo})

        return "\n\n".join(trechos), fontes
    except Exception as e:
        return "(Erro ao buscar no RAG: " + str(e) + ")", []


# ======================
# Excel Matriz_FP – Lookup determinístico por PN
# ======================
@st.cache_data(show_spinner="Carregando Matriz FP...")
def carregar_df_fp_matriz(file_path: str, sheet_name: str) -> "pd.DataFrame":
    """
    Carrega a matriz FP. Aceita CSV ou Excel.
    Prioriza CSV se existir no mesmo diretório com mesmo nome base.
    """
    # Colunas esperadas na planilha (exatas 35 colunas)
    _COLUNAS_CSV = [
        "marketing_project","description","pn_gen","pn_alternative","pn_nxp","pn_fleetpro",
        "launched_quarter","launched_month","ano","modal","ag","ce","am",
        "tractor_case_ih","combine_case_ih","headers_case_ih","sch_case_ih",
        "sprayers_case_ih","planters_case_ih","other_machines_case_ih",
        "tractor_nhag","combine_nhag","headers_nhag","sprayers_nhag","planters_nhag",
        "forage_balers_and_others_nhag","other_machines_nhag",
        "montadora","john_deere","macdon","agco","ideal","massey_ferguson","valtra","cd_tech_type",
    ]

    # Verifica se existe versão CSV (menor e mais rápida)
    csv_path = os.path.splitext(file_path)[0] + ".csv"
    if os.path.exists(csv_path):
        try:
            # Lê apenas as primeiras 35 colunas, ignorando colunas extras vazias do final
            df = pd.read_csv(
                csv_path,
                encoding="utf-8",
                low_memory=False,
                usecols=range(len(_COLUNAS_CSV)),   # ignora colunas extras
                names=_COLUNAS_CSV,                  # força nomes corretos
                header=0,                            # primeira linha é cabeçalho original
                on_bad_lines="skip",
            )
        except UnicodeDecodeError:
            df = pd.read_csv(
                csv_path,
                encoding="latin1",
                low_memory=False,
                usecols=range(len(_COLUNAS_CSV)),
                names=_COLUNAS_CSV,
                header=0,
                on_bad_lines="skip",
            )
    else:
        df = pd.read_excel(file_path, sheet_name=sheet_name, engine="openpyxl")

    # Normaliza nomes de colunas: lowercase + underscore
    df.columns = [
        str(c).strip().lower()
        .replace(" ", "_").replace("-", "_").replace("/", "_").replace(".", "_")
        for c in df.columns
    ]
    for c in COLUNAS_BUSCA_PN:
        if c not in df.columns:
            df[c] = ""
    return df


# Colunas prioritárias para exibir no resultado
_COLUNAS_EXIBIR = [
    "description", "pn_fleetpro", "pn_gen", "pn_alternative", "pn_nxp",
    "marketing_project", "montadora", "modal", "ag", "ce", "am",
    "cd_tech_type", "launched_quarter", "launched_month", "ano",
    "tractor_case_ih", "combine_case_ih", "headers_case_ih",
    "sch_case_ih", "sprayers_case_ih", "planters_case_ih", "other_machines_case_ih",
    "tractor_nhag", "combine_nhag", "headers_nhag",
    "sprayers_nhag", "planters_nhag", "forage_balers_and_others_nhag", "other_machines_nhag",
    "john_deere", "macdon", "agco", "ideal", "massey_ferguson", "valtra",
]

def formatar_linha_como_lista(df: "pd.DataFrame", row_index_df: int) -> str:
    if row_index_df < 0 or row_index_df >= len(df):
        return "Índice fora do intervalo."
    row = df.iloc[row_index_df]
    partes = []
    # Colunas que contêm listas de modelos separados por ;
    _COLUNAS_LISTA = set(COLUNAS_MODELOS_EQUIP) | {"cd_tech_type"}

    # Exibe colunas prioritárias primeiro, depois o restante
    cols_usadas = []
    for col in _COLUNAS_EXIBIR:
        if col in df.columns:
            val = row[col]
            if not pd.isna(val) and str(val).strip() not in ("", "-", "nan", "NaN"):
                val_str = _formatar_lista_equip(str(val)) if col in _COLUNAS_LISTA else str(val).strip()
                col_label = col.replace("_", " ").title()
                partes.append(f"- **{col_label}**: {val_str}")
            cols_usadas.append(col)
    # Demais colunas não listadas acima
    for col in df.columns:
        if col not in cols_usadas:
            val = row[col]
            if not pd.isna(val) and str(val).strip() not in ("", "-", "nan", "NaN"):
                partes.append(f"- {col}: {str(val).strip()}")
    return "\n".join(partes) if partes else "Sem dados disponíveis."


def _formatar_lista_equip(val: str) -> str:
    """Formata lista de equipamentos separada por ; em formato legível."""
    if not val or str(val).strip() in ("", "-", "nan", "NaN"):
        return ""
    itens = [x.strip() for x in str(val).split(";") if x.strip()]
    if len(itens) <= 1:
        return val
    return ", ".join(itens)


def detectar_busca_por_equipamento(mensagem: str):
    msg_upper = mensagem.upper()
    chaves_ordenadas = sorted(MAPA_PALAVRAS_COLUNAS.keys(), key=lambda x: -len(x))
    colunas_encontradas = []
    for chave in chaves_ordenadas:
        if chave in msg_upper:
            for col in MAPA_PALAVRAS_COLUNAS[chave]:
                if col not in colunas_encontradas:
                    colunas_encontradas.append(col)
    return colunas_encontradas


def detectar_busca_marketing(mensagem: str):
    """
    Detecta se a mensagem é uma consulta por categoria/tipo de produto.
    Normaliza acentos para comparação e retorna o termo EXATO da coluna
    MARKETING PROJECT (ou o sinônimo detectado), ou None se não encontrar.

    Exemplos:
        "quais rolamentos FleetPro vocês têm?"    → "ROLAMENTOS"
        "preciso de correias para minha colheitadeira" → "CORREIAS"
        "tem filtro disponível?"                  → "FILTROS"
        "undercarriage disponível?"               → "UNDERCARRIAGE"
    """
    import unicodedata

    def normalizar(s: str) -> str:
        """Remove acentos e converte para maiúsculo."""
        return unicodedata.normalize("NFD", s.upper()).encode("ascii", "ignore").decode("ascii")

    msg_norm = normalizar(mensagem)

    # Ordena: mais longo primeiro (evita "BELT" antes de "BELTS", "CORREIA" antes de "CORREIAS")
    termos_ordenados = sorted(PALAVRAS_MARKETING, key=lambda x: -len(x))

    for termo in termos_ordenados:
        termo_norm = normalizar(termo)
        # Usa \b apenas para termos puramente alfanuméricos; para termos com & ou espaço, usa in
        if re.search(r"[^A-Z0-9]", termo_norm):
            # Termo contém espaço, & ou outro caractere especial → busca direta
            if termo_norm in msg_norm:
                return termo
        else:
            # Termo simples → busca como palavra inteira
            if re.search(r"\b" + re.escape(termo_norm) + r"\b", msg_norm):
                return termo

    return None


def _extrair_palavras_extras(mensagem: str, termo_mkt: str) -> list:
    """
    Extrai tokens da mensagem que não são o termo de marketing nem palavras genéricas.
    Permite refinar a busca na DESCRIPTION (ex: "rolamento 6205" → ["6205"]).
    """
    stopwords_busca = {
        "QUAIS", "QUAL", "PN", "PNS", "PARA", "NO", "NA", "OS", "AS", "DE",
        "DO", "DA", "UM", "UMA", "QUE", "SAO", "SÃO", "ME", "MOSTRE", "MOSTRA",
        "LISTA", "LISTAR", "VER", "TENHO", "PRECISO", "QUERO", "TEM", "FLEETPRO",
        "DISPONIVEIS", "DISPONÍVEIS", "VOCÊS", "VOCES", "TEMOS", "HAI", "HÁ",
        "EXISTE", "EXISTEM", "BUSCA", "BUSCAR", "ENCONTRA", "ACHAR", "THE", "AND",
        "FOR", "WITH", "HAVE", "SHOW", "FIND", "GET", "ALL",
    } | {termo_mkt}
    tokens = re.findall(r"[A-Za-z0-9]+", mensagem.upper())
    extras = [
        t for t in tokens
        if t not in stopwords_busca
        and len(t) >= 3
        and t not in {p.upper() for p in PALAVRAS_MARKETING}
    ]
    return list(dict.fromkeys(extras))




def buscar_equip_e_marketing(
    df: "pd.DataFrame",
    mensagem: str,
    colunas_equip: list,
    termo_marketing: str,
    max_resultados: int = 50,
) -> str:
    """
    Busca combinada: filtra por categoria de produto (marketing_project)
    E por equipamento (colunas de modelo). Ex: rolamento + trator case ih.
    """
    import unicodedata

    def norm_sem_acento(s: str) -> str:
        return unicodedata.normalize("NFD", s.upper()).encode("ascii", "ignore").decode("ascii")

    # Resolve sinônimo de marketing
    termo_upper = termo_marketing.upper()
    termo_busca = MAPA_SINONIMOS_MARKETING.get(termo_upper, termo_upper)
    if termo_busca == termo_upper:
        termo_busca = MAPA_SINONIMOS_MARKETING.get(norm_sem_acento(termo_upper), termo_upper)

    # Filtro 1: categoria (marketing_project)
    if COLUNA_MARKETING not in df.columns:
        return buscar_por_equipamento(df, mensagem, colunas_equip, max_resultados)

    mask_mkt = (
        df[COLUNA_MARKETING].notna()
        & df[COLUNA_MARKETING].astype(str).str.upper().str.contains(
            re.escape(termo_busca), na=False
        )
    )
    df_filtrado = df[mask_mkt].copy()

    if df_filtrado.empty:
        # Fallback: só por equipamento
        return buscar_por_equipamento(df, mensagem, colunas_equip, max_resultados)

    # Filtro 2: equipamento (dentro do df já filtrado por categoria)
    # Extrai modelo da mensagem para filtrar
    tokens = re.findall(r"[A-Za-z0-9]+", mensagem.upper())
    palavras_genericas = {
        "QUAIS", "QUAL", "PN", "PARA", "DE", "DO", "DA", "UM", "UMA",
        "ROLAMENTO", "ROLAMENTOS", "FILTRO", "FILTROS", "CORREIA", "CORREIAS",
        "CASE", "NHAG", "TRATOR", "TRACTOR", "COMBINE", "COLHEITADEIRA",
        "NEW", "HOLLAND", "IH", "PRECISO", "QUERO", "TEM", "HAI",
    }
    candidatos_modelo = [t for t in tokens if t not in palavras_genericas and len(t) >= 3]

    mask_equip = pd.Series([False] * len(df_filtrado), index=df_filtrado.index)
    for col in colunas_equip:
        if col in df_filtrado.columns:
            mask_col = df_filtrado[col].notna() & ~df_filtrado[col].isin(["-", "", "NaN", "nan"])
            if candidatos_modelo and col in COLUNAS_MODELOS_EQUIP:
                def contem_modelo_combo(val):
                    if pd.isna(val): return False
                    itens = [x.strip().upper() for x in str(val).split(";")]
                    return any(any(cm in item for cm in candidatos_modelo) for item in itens)
                mask_col = mask_col & df_filtrado[col].apply(contem_modelo_combo)
            mask_equip = mask_equip | mask_col

    df_combinado = df_filtrado[mask_equip].copy()

    # Se filtro combinado vazio, usa só o de marketing
    if df_combinado.empty:
        df_combinado = df_filtrado.copy()

    if len(df_combinado) > max_resultados:
        df_combinado = df_combinado.head(max_resultados)

    total = len(df_combinado)
    cols_label = ", ".join(c.replace("_", " ").upper() for c in colunas_equip[:2])
    partes = [f"Encontrei **{total}** produto(s) — **{termo_busca}** para **{cols_label}**:\n"]

    for i, (idx_df, _) in enumerate(df_combinado.iterrows(), start=1):
        row_index_df = df.index.get_loc(idx_df)
        partes.append(f"### Produto {i}\n")
        partes.append(formatar_linha_como_lista(df, int(row_index_df)))
        partes.append("")

    return "\n".join(partes)

def _normalizar_tt(x) -> str:
    """Normaliza TT code: remove espaços e converte para string limpa."""
    if x is None:
        return ""
    return re.sub(r"\s+", "", str(x).strip())


def buscar_por_tt_code(df: "pd.DataFrame", tt_code: str, max_resultados: int = 50) -> str:
    """
    Busca por TT code (cd_tech_type).
    Cada célula pode conter múltiplos códigos separados por vírgula ou ponto-e-vírgula.
    """
    if "cd_tech_type" not in df.columns:
        return "⚠️ Coluna `cd_tech_type` não encontrada na planilha."

    tt_norm = _normalizar_tt(tt_code)
    if not tt_norm:
        return ""

    # Expande células com múltiplos TT codes e verifica match
    def contem_tt(val):
        if pd.isna(val):
            return False
        # Split por vírgula, ponto-e-vírgula ou espaço
        codigos = re.split(r"[,;]+", str(val).strip())
        return any(_normalizar_tt(c.strip()) == tt_norm for c in codigos if c.strip())

    mask = df["cd_tech_type"].apply(contem_tt)
    encontrados = df[mask]

    if encontrados.empty:
        return f"Não encontrei nenhum produto para o TT code `{tt_code}` na matriz FP."

    if len(encontrados) > max_resultados:
        encontrados = encontrados.head(max_resultados)

    partes = [f"Encontrei **{len(encontrados)}** produto(s) para o TT code `{tt_code}`:\n"]
    for i, (idx_df, _) in enumerate(encontrados.iterrows(), start=1):
        row_index_df = df.index.get_loc(idx_df)
        partes.append(f"### Produto {i}\n")
        partes.append(formatar_linha_como_lista(df, int(row_index_df)))
        partes.append("")
    return "\n".join(partes)


def _detectar_pergunta_contagem(mensagem: str) -> bool:
    palavras = ["quantos", "quantidade", "total de itens", "total de produtos",
                "portfólio", "portfolio", "quantas peças", "quantas pecas",
                "tamanho do portfólio", "itens no portfólio", "itens temos",
                "peças temos", "pecas temos", "produtos temos",
                "tem no portfólio", "tem no portfolio", "quantos itens tem",
                "quantos produtos tem", "itens fleetpro", "produtos fleetpro tem",
                "itens cadastrados", "produtos cadastrados", "peças fleetpro tem",
                "pecas fleetpro tem", "quantos produtos", "quantas pecas tem",
                "quantas peças tem", "quantos itens", "total de peças",
                "total de pecas", "numero de itens", "número de itens"]
    msg = mensagem.lower()
    return any(p in msg for p in palavras)


def detectar_tt_code(mensagem: str) -> str:
    """
    Detecta TT code na mensagem — sequência numérica de 4+ dígitos
    precedida por 'TT', 'tt', 'tech type', 'techtype' ou isolada como número puro.
    Retorna o código detectado ou string vazia.
    """
    # Padrão explícito: TT seguido de número
    m = re.search(r"\b[Tt][Tt]\s*(\d{4,})", mensagem)
    if m:
        return m.group(1)
    # "tech type 12345" ou "techtype 12345"
    m = re.search(r"tech\s*type\s*(\d{4,})", mensagem, re.IGNORECASE)
    if m:
        return m.group(1)
    # Número puro de 5+ dígitos que não seja PN (PNs têm letras)
    # Só detecta se a mensagem mencionar "TT" ou "tech"
    if re.search(r"\b(tt|tech)\b", mensagem, re.IGNORECASE):
        m = re.search(r"\b(\d{4,})\b", mensagem)
        if m:
            return m.group(1)
    return ""

def buscar_por_marketing(
    df: "pd.DataFrame",
    termo: str,
    mensagem: str,
    max_resultados: int = 100,
) -> str:
    """
    Busca linhas onde MARKETING PROJECT contém o termo (parcial, case-insensitive).
    Resolve sinônimos para o valor exato da coluna antes de buscar.
    Aplica filtro extra na DESCRIPTION se o usuário informou termos adicionais.
    Retorna markdown formatado.
    """
    import unicodedata

    def norm_sem_acento(s: str) -> str:
        return unicodedata.normalize("NFD", s.upper()).encode("ascii", "ignore").decode("ascii")

    # Resolve sinônimo → valor exato da coluna (ex: "ROLAMENTO" → "ROLAMENTOS")
    termo_upper = termo.upper()
    termo_busca = MAPA_SINONIMOS_MARKETING.get(termo_upper, termo_upper)
    # Tenta também sem acento como fallback
    if termo_busca == termo_upper:
        termo_busca = MAPA_SINONIMOS_MARKETING.get(norm_sem_acento(termo_upper), termo_upper)

    if COLUNA_MARKETING not in df.columns:
        return (
            f"⚠️ A coluna `{COLUNA_MARKETING}` não foi encontrada na planilha. "
            "Verifique se o nome está correto."
        )

    mask_mkt = (
        df[COLUNA_MARKETING].notna()
        & df[COLUNA_MARKETING].astype(str).str.upper().str.contains(
            re.escape(termo_busca), na=False
        )
    )
    df_filtrado = df[mask_mkt].copy()

    if df_filtrado.empty:
        # Tenta busca parcial sem acento como último recurso
        mask_fallback = (
            df[COLUNA_MARKETING].notna()
            & df[COLUNA_MARKETING].astype(str).apply(norm_sem_acento).str.contains(
                re.escape(norm_sem_acento(termo_busca)), na=False
            )
        )
        df_filtrado = df[mask_fallback].copy()

    if df_filtrado.empty:
        return (
            f"Não encontrei nenhum produto na categoria **{termo_busca}** na matriz FP.\n\n"
            f"Categorias disponíveis: {', '.join(f'`{v}`' for v in sorted(set(MAPA_SINONIMOS_MARKETING.values())))}"
        )

    # Filtro secundário: termos extras na DESCRIPTION
    palavras_extra = _extrair_palavras_extras(mensagem, termo)
    aplicou_filtro_extra = False
    if palavras_extra and "DESCRIPTION" in df_filtrado.columns:
        pattern_extra = "|".join(re.escape(p) for p in palavras_extra)
        mask_desc = df_filtrado["DESCRIPTION"].astype(str).str.upper().str.contains(
            pattern_extra, na=False
        )
        if mask_desc.any():
            df_filtrado = df_filtrado[mask_desc]
            aplicou_filtro_extra = True

    truncado = False
    if len(df_filtrado) > max_resultados:
        df_filtrado = df_filtrado.head(max_resultados)
        truncado = True

    total = len(df_filtrado)
    sufixo_trunc = f" *(limitado aos primeiros {max_resultados})*" if truncado else ""
    subtitulo_extra = (
        f" — filtro adicional: `{'`, `'.join(palavras_extra)}`"
        if aplicou_filtro_extra
        else ""
    )

    linhas = [
        f"## 🔍 {COLUNA_MARKETING}: **{termo_busca}** — {total} produto(s){subtitulo_extra}{sufixo_trunc}\n"
    ]

    for _, row in df_filtrado.iterrows():
        pn_fp   = str(row.get("pn_fleetpro", "") or "").strip()
        pn_gen  = str(row.get("pn_gen", "")  or "").strip()
        descr   = str(row.get("description", "") or "").strip()
        mkt_val = str(row.get("marketing_project", "") or "").strip()

        partes = [f"- **{descr}**" if descr else "- *(sem descrição)*"]
        if mkt_val and mkt_val not in ("-", ""):
            partes.append(f"  Categoria: `{mkt_val}`")
        if pn_fp and pn_fp not in ("-", ""):
            partes.append(f"  PN FleetPro: `{pn_fp}`")
        if pn_gen and pn_gen not in ("-", ""):
            partes.append(f"  PN GEN: `{pn_gen}`")

        linhas.append("\n".join(partes))

    linhas.append("")
    return "\n".join(linhas)


def buscar_por_equipamento(df: "pd.DataFrame", mensagem: str, colunas_alvo: list, max_resultados: int = 100) -> str:
    import re as _re

    modelo_filtro = None
    tokens = _re.findall(r'[A-Za-z0-9]+', mensagem.upper())
    palavras_genericas = {
        'QUAIS', 'QUAL', 'PN', 'PNS', 'PARA', 'NO', 'NA', 'OS', 'AS', 'DE',
        'DO', 'DA', 'UM', 'UMA', 'QUE', 'SAO', 'SÃO', 'POSSO', 'USAR',
        'UTILIZAR', 'LISTA', 'LISTAR', 'MOSTRAR', 'MOSTRAME', 'QUERO',
        'VER', 'EQUIPAMENTO', 'MAQUINA', 'JOHN', 'DEERE', 'CASE', 'NHAG',
        'MASSEY', 'FERGUSON', 'VALTRA', 'AGCO', 'MACDON', 'IDEAL',
        'TRATOR', 'TRACTOR', 'COMBINE', 'COLHEITADEIRA', 'HEADER',
        'PLATAFORMA', 'PULVERIZADOR', 'SPRAYER', 'PLANTADEIRA', 'PLANTER',
        'NEW', 'HOLLAND', 'COMPATIVEIS', 'COMPATÍVEIS', 'IH', 'ME', 'MOSTRE',
    }
    candidatos_modelo = [t for t in tokens if t not in palavras_genericas and len(t) >= 3]
    if candidatos_modelo:
        modelo_filtro = max(candidatos_modelo, key=len)

    resultados_por_coluna = {}

    for col in colunas_alvo:
        if col not in df.columns:
            continue

        mask = df[col].notna() & ~df[col].isin(['-', '', 'NaN', 'nan'])
        df_col = df[mask].copy()

        if df_col.empty:
            continue

        if modelo_filtro and col in COLUNAS_MODELOS_EQUIP:
            # Busca em listas separadas por ponto-e-vírgula dentro da célula
            def contem_modelo(val):
                if pd.isna(val):
                    return False
                itens = [x.strip().upper() for x in str(val).split(";")]
                return any(modelo_filtro in item for item in itens)
            mask_modelo = df_col[col].apply(contem_modelo)
            if mask_modelo.any():
                df_col = df_col[mask_modelo]

        if len(df_col) > max_resultados:
            df_col = df_col.head(max_resultados)
            truncado = True
        else:
            truncado = False

        resultados_por_coluna[col] = (df_col, truncado)

    if not resultados_por_coluna:
        return (
            "Não encontrei peças para esse equipamento/modelo na Matriz FP. "
            "Tente especificar o equipamento completo, por exemplo: "
            "'COMBINE - CASE IH 2188', 'TRACTOR - CASE IH PUMA 150', "
            "'COMBINE - NEW HOLLAND'. Você também pode buscar por PN ou categoria de produto."
        )

    linhas_output = []
    total_itens = 0

    for col, (df_resultado, truncado) in resultados_por_coluna.items():
        e_coluna_marca = col in COLUNAS_MARCAS_PN
        sufixo = f"  *(limitado aos primeiros {max_resultados})*" if truncado else ""
        col_label = col.replace("_", " ").upper()
        linhas_output.append(f"## 🔧 {col_label} — {len(df_resultado)} peça(s){sufixo}\n")

        for _, row in df_resultado.iterrows():
            pn_gen = str(row.get('pn_gen', '') or '').strip()
            pn_fp = str(row.get('pn_fleetpro', '') or '').strip()
            descricao = str(row.get('description', '') or '').strip()
            val_col = str(row.get(col, '') or '').strip()

            partes_linha = [f"- **{descricao}**"]
            if pn_fp and pn_fp not in ('-', ''):
                partes_linha.append(f"  PN FleetPro: `{pn_fp}`")
            if pn_gen and pn_gen not in ('-', ''):
                partes_linha.append(f"  PN GEN: `{pn_gen}`")
            if e_coluna_marca and val_col and val_col not in ('-', ''):
                partes_linha.append(f"  PN {col}: `{val_col}`")
            elif not e_coluna_marca and val_col and val_col not in ('-', ''):
                partes_linha.append(f"  Modelos: {val_col}")

            linhas_output.append("\n".join(partes_linha))
            total_itens += 1

        linhas_output.append("")

    header = f"Encontrei **{total_itens}** peça(s) compatível(is):\n"
    if modelo_filtro and any(col in COLUNAS_MODELOS_EQUIP for col in colunas_alvo):
        header = f"Encontrei **{total_itens}** peça(s) compatível(is) (filtro de modelo: `{modelo_filtro}`):\n"

    return header + "\n".join(linhas_output)


def buscar_por_description(df: "pd.DataFrame", mensagem: str, max_resultados: int = 30) -> str:
    """
    Busca fallback por texto livre na coluna DESCRIPTION.
    Retorna string formatada ou vazia se não encontrar nada relevante.
    """
    import unicodedata

    def norm(s: str) -> str:
        return unicodedata.normalize("NFD", s.upper()).encode("ascii", "ignore").decode("ascii")

    stopwords = {
        "QUAIS", "QUAL", "PN", "PNS", "PARA", "NO", "NA", "OS", "AS", "DE",
        "DO", "DA", "UM", "UMA", "QUE", "SAO", "ME", "MOSTRE", "MOSTRA",
        "LISTA", "LISTAR", "VER", "TENHO", "PRECISO", "QUERO", "TEM", "FLEETPRO",
        "DISPONIVEIS", "DISPONIVEL", "VOCES", "TEMOS", "HAI", "EXISTE", "EXISTEM",
        "BUSCA", "BUSCAR", "ENCONTRA", "ACHAR", "THE", "AND", "FOR", "WITH",
        "HAVE", "SHOW", "FIND", "GET", "ALL", "CASO", "ESSE", "ESSA", "ISSO",
    }

    # Mapeamento de sinônimos: variação do usuário → termo na base
    _SINONIMOS = {
        "SILOBOLSA": "SILOBAG",
        "SILO BOLSA": "SILOBAG",
        "BOLSA SILO": "SILOBAG",
        "BOLSA": "SILOBAG",
    }

    # Normaliza a mensagem aplicando sinônimos antes de tokenizar
    msg_norm = norm(mensagem)
    for variacao, substituto in _SINONIMOS.items():
        msg_norm = msg_norm.replace(norm(variacao), substituto)

    tokens = [t for t in re.findall(r"[A-Za-z0-9]+", msg_norm)
              if norm(t) not in stopwords and len(t) >= 3]

    if not tokens or "description" not in df.columns:
        return ""

    # Filtra linhas onde DESCRIPTION contém algum dos tokens
    desc_upper = df["description"].fillna("").astype(str).apply(norm)
    mask = pd.Series([False] * len(df), index=df.index)
    for token in tokens:
        mask = mask | desc_upper.str.contains(norm(token), na=False)

    df_match = df[mask].head(max_resultados)
    if df_match.empty:
        return ""

    linhas = [f"Encontrei **{len(df_match)}** item(ns) com descrição relacionada:\n"]
    for _, row in df_match.iterrows():
        descricao = str(row.get("description", "") or "").strip()
        pn_fp = str(row.get("pn_fleetpro", "") or "").strip()
        pn_gen = str(row.get("pn_gen", "") or "").strip()
        mkt = str(row.get("marketing_project", "") or "").strip()

        partes = [f"- **{descricao}**"]
        if pn_fp and pn_fp not in ("-", ""):
            partes.append(f"  PN FleetPro: `{pn_fp}`")
        if pn_gen and pn_gen not in ("-", ""):
            partes.append(f"  PN GEN: `{pn_gen}`")
        if mkt and mkt not in ("-", ""):
            partes.append(f"  Categoria: {mkt}")
        linhas.append("\n".join(partes))

    return "\n".join(linhas)


def _buscar_pn_no_df(df: "pd.DataFrame", pn_norm: str, max_resultados: int) -> "pd.DataFrame":
    df_busca = df[COLUNAS_BUSCA_PN].copy()
    for c in COLUNAS_BUSCA_PN:
        df_busca[c] = df_busca[c].apply(norm_pn)

    mask = False
    for c in COLUNAS_BUSCA_PN:
        mask = mask | (df_busca[c] == pn_norm)

    encontrados = df[mask]
    if len(encontrados) > max_resultados:
        encontrados = encontrados.head(max_resultados)
    return encontrados


def _formatar_resultados(df: "pd.DataFrame", encontrados: "pd.DataFrame", pn_exibido: str) -> str:
    partes = [f"Encontrei **{len(encontrados)}** opção(ões) para o PN `{pn_exibido}`:\n"]
    for i, (idx_df, _row) in enumerate(encontrados.iterrows(), start=1):
        row_index_df = df.index.get_loc(idx_df)
        partes.append(f"### Opção {i}\n")
        partes.append(formatar_linha_como_lista(df, int(row_index_df)))
        partes.append("")
    return "\n".join(partes)


def procurar_pn(df: "pd.DataFrame", mensagem_usuario: str, max_resultados: int = 50) -> str:
    if not mensagem_usuario or not mensagem_usuario.strip():
        return "Informe um PN válido para pesquisa."

    pn_direto = norm_pn(mensagem_usuario)
    if pn_direto:
        encontrados = _buscar_pn_no_df(df, pn_direto, max_resultados)
        if not encontrados.empty:
            return _formatar_resultados(df, encontrados, mensagem_usuario.strip())

    candidatos = extrair_pns_da_mensagem(mensagem_usuario)

    for candidato in candidatos:
        encontrados = _buscar_pn_no_df(df, candidato, max_resultados)
        if not encontrados.empty:
            return _formatar_resultados(df, encontrados, candidato)

    if candidatos:
        tentados = ", ".join(f"`{c}`" for c in candidatos[:5])
        return (
            f"Não encontrei nenhum dos seguintes PNs na matriz: {tentados}.\n\n"
            "Verifique se o número está correto ou tente digitar apenas o código da peça."
        )

    return ""


# ======================
# Inicializar LLM
# ======================
def inicializar_FleetPro(provedor: str, modelo: str, api_key: str):
    api_key = (api_key or "").strip()
    if not api_key:
        st.error("Informe a API key antes de inicializar.")
        return

    chat_cls = CONFIG_MODELOS[provedor]["chat"]
    chat = chat_cls(model=modelo, api_key=api_key)

    st.session_state["chat"] = chat
    st.session_state["provedor"] = provedor
    st.session_state["modelo"] = modelo
    st.session_state["api_key_openai_rag"] = api_key if provedor == "OpenAI" else st.session_state.get("api_key_openai_rag", "")
    st.session_state.setdefault("memoria", ConversationBufferMemory())

    if "session_id" not in st.session_state:
        import uuid, datetime as _dt
        st.session_state["session_id"] = str(uuid.uuid4())[:8]
        st.session_state["session_inicio"] = _dt.datetime.now().isoformat()
        st.session_state["tokens_sessao"] = 0
        st.session_state["n_mensagens_sessao"] = 0

    resetar_vectorstore()
    with st.spinner("Indexando documentos..."):
        vs = obter_vectorstore()

    if vs:
        st.success(f"✅ Agente inicializado: {provedor} / {modelo} — documentos indexados.")
    else:
        st.warning(f"✅ Agente inicializado: {provedor} / {modelo} — RAG indisponível (sem documentos ou rede bloqueada).")


# ======================
# Detecção de busca simples (sem intenção de argumentação)
# ======================
_VERBOS_ARGUMENTACAO = {
    "beneficio", "benefícios", "vantagem", "vantagens", "argumento", "argumentos",
    "por que", "porque", "convencer", "vender", "cross", "objecao", "objeção",
    "diferencial", "diferenciais", "qualidade", "especificacao", "especificações",
    "comparar", "versus", "melhor", "recomend", "indicar", "sugerir",
}

def _eh_busca_simples(mensagem: str) -> bool:
    """
    Retorna True se a mensagem é uma busca direta por produto/PN,
    sem intenção de receber argumentação ou análise do LLM.
    """
    msg_lower = mensagem.lower().strip()
    # Se contém verbo de argumentação, não é busca simples
    for verbo in _VERBOS_ARGUMENTACAO:
        if verbo in msg_lower:
            return False
    # Se tem ponto de interrogação, é uma pergunta → LLM responde
    if "?" in msg_lower:
        return False
    # Mensagem curta sem verbo → busca simples (ex: "silo bolsa", "1890DSH", "rolamento")
    palavras = msg_lower.split()
    if len(palavras) <= 4:
        return True
    return False


# ======================
# UI – Chat principal
# ======================
def pagina_chat():
    col_logo, col_titulo = st.columns([1, 6])
    with col_logo:
        st.image("base_docs/fleetpro_logo.png", width=120)
    with col_titulo:
        st.header("🕵️‍♂️ FleetPro Expert 🛠️", divider=True)

    chat_model = st.session_state.get("chat")
    memoria: ConversationBufferMemory = st.session_state.get("memoria", ConversationBufferMemory())

    for mensagem in memoria.buffer_as_messages:
        st.chat_message(mensagem.type).markdown(mensagem.content)

    # ── Argumentação de venda pendente (clique no botão) ─────────────────
    if st.session_state.get("_argumentacao_pendente"):
        _arg = st.session_state.pop("_argumentacao_pendente")
        _input_arg = _arg["input_usuario"]
        _resultado_arg = _arg["resultado_matriz"]
        _contexto_arg = _arg.get("contexto_rag", "")
        chat_model_arg = st.session_state.get("chat")
        if chat_model_arg is not None:
            _blocos_arg = ["## Resultado da Busca na Matriz FP\n\n" + _resultado_arg]
            if _contexto_arg:
                _blocos_arg.append("## Conhecimento Adicional\n\n" + _contexto_arg)
            _ctx_arg = "\n\n---\n\n".join(_blocos_arg)
            _prompt_arg = (
                f"Você é o FleetPro Expert, assistente de vendas para vendedores de balcão.\n\n"
                f"O vendedor quer argumentos de venda para o produto: **{_input_arg}**\n\n"
                f"CONTEXTO SOBRE FLEETPRO:\n"
                f"- FleetPro é a marca própria CNH Industrial para máquinas Case IH e New Holland.\n"
                f"- VV, TEEJET, IPESA, REXNORD, PETRONAS são FORNECEDORES — sempre apresente como produtos FleetPro.\n\n"
                f"INSTRUÇÕES:\n"
                f"- Use APENAS os dados abaixo. Não invente informações ausentes.\n"
                f"- Estruture a resposta em seções curtas com bullet points:\n"
                f"  1. **Opções disponíveis** — liste os PNs e descrições de forma limpa\n"
                f"  2. **Argumentos de venda** — 3 a 5 bullets diretos para usar no balcão\n"
                f"  3. **Como rebater objeções** — 2 bullets (preço / qualidade)\n"
                f"  4. **Cross-selling** — 1 ou 2 produtos complementares se houver na base\n"
                f"- Seja direto. O vendedor está atendendo o cliente agora.\n\n"
                f"Dados das fontes:\n\n{_ctx_arg}"
            )
            with st.chat_message("ai"):
                import time as _time
                _t0 = _time.time()
                _resp_arg = chat_model_arg.invoke(_prompt_arg)
                _tempo_arg = round(_time.time() - _t0, 2)
                _resposta_arg = _resp_arg.content
                st.markdown(_resposta_arg)
            memoria.chat_memory.add_user_message(f"[Argumentação] {_input_arg}")
            memoria.chat_memory.add_ai_message(_resposta_arg)
            st.session_state["memoria"] = memoria
            _uso_arg = getattr(_resp_arg, "usage_metadata", {}) or {}
            _in_arg = _uso_arg.get("input_tokens", 0)
            _out_arg = _uso_arg.get("output_tokens", 0)
            _tot_arg = _uso_arg.get("total_tokens", _in_arg + _out_arg)
            st.session_state["tokens_sessao"] = st.session_state.get("tokens_sessao", 0) + _tot_arg
            st.session_state["n_mensagens_sessao"] = st.session_state.get("n_mensagens_sessao", 0) + 1
            import threading as _threading
            _threading.Thread(
                target=_salvar_log_tokens,
                args=({
                    "session_id": st.session_state.get("session_id", ""),
                    "timestamp_inicio_sessao": st.session_state.get("session_inicio", ""),
                    "timestamp_resposta": _time.strftime("%Y-%m-%dT%H:%M:%S"),
                    "modelo": st.session_state.get("modelo", ""),
                    "provedor": st.session_state.get("provedor", ""),
                    "perfil_usuario": st.session_state.get("perfil_usuario", ""),
                    "input_tokens": _in_arg,
                    "output_tokens": _out_arg,
                    "total_tokens": _tot_arg,
                    "tempo_resposta_seg": _tempo_arg,
                    "tokens_acumulados_sessao": st.session_state.get("tokens_sessao", 0),
                    "n_mensagens_sessao": st.session_state.get("n_mensagens_sessao", 0),
                },),
                daemon=True,
            ).start()

    # ── Botões de perfil (exibe apenas enquanto perfil não definido) ──────
    if st.session_state.get("perfil_usuario") is None:
        st.caption("Quem é você?")
        col_v, col_c, _ = st.columns([1, 1, 4])
        with col_v:
            if st.button("🧑‍💼 Vendedor", use_container_width=True, key="btn_perfil_vendedor"):
                st.session_state["perfil_usuario"] = "vendedor"
                resposta = "✅ Modo **Vendedor de Balcão** ativado. Pode fazer sua pergunta!"
                st.chat_message("human").markdown("Vendedor")
                st.chat_message("ai").markdown(resposta)
                memoria.chat_memory.add_user_message("Vendedor")
                memoria.chat_memory.add_ai_message(resposta)
                st.session_state["memoria"] = memoria
                st.rerun()
        with col_c:
            if st.button("🌾 Cliente", use_container_width=True, key="btn_perfil_cliente"):
                st.session_state["perfil_usuario"] = "usuario"
                resposta = "✅ Modo **Cliente / Usuário da Peça** ativado. Pode fazer sua pergunta!"
                st.chat_message("human").markdown("Cliente")
                st.chat_message("ai").markdown(resposta)
                memoria.chat_memory.add_user_message("Cliente")
                memoria.chat_memory.add_ai_message(resposta)
                st.session_state["memoria"] = memoria
                st.rerun()

    input_usuario = st.chat_input(
        "Digite um PN para busca, ou uma pergunta sobre produtos, objeções e recomendações..."
    )

    if not input_usuario:
        return

    # ── Captura escolha de perfil via texto (fallback) ────────────────────
    if st.session_state.get("perfil_usuario") is None:
        resposta_lower = input_usuario.lower().strip()
        if any(p in resposta_lower for p in ["1", "vendedor", "vendo", "balcão", "balcao", "revend"]):
            st.session_state["perfil_usuario"] = "vendedor"
            st.chat_message("human").markdown(input_usuario)
            with st.chat_message("ai"):
                resposta = "✅ Perfeito! Modo **Vendedor de Balcão** ativado. Pode fazer sua pergunta!"
                st.markdown(resposta)
            memoria.chat_memory.add_user_message(input_usuario)
            memoria.chat_memory.add_ai_message(resposta)
            st.session_state["memoria"] = memoria
            st.stop()
        elif any(p in resposta_lower for p in ["2", "usuario", "usuário", "uso", "minha máquina", "minha maquina", "agricultor", "operador", "cliente"]):
            st.session_state["perfil_usuario"] = "usuario"
            st.chat_message("human").markdown(input_usuario)
            with st.chat_message("ai"):
                resposta = "✅ Perfeito! Modo **Cliente / Usuário da Peça** ativado. Pode fazer sua pergunta!"
                st.markdown(resposta)
            memoria.chat_memory.add_user_message(input_usuario)
            memoria.chat_memory.add_ai_message(resposta)
            st.session_state["memoria"] = memoria
            st.stop()

    st.chat_message("human").markdown(input_usuario)

    with st.chat_message("ai"):
        try:
            usar_fp_matriz = st.session_state.get("usar_fp_matriz", True)
            usar_rag = st.session_state.get("usar_rag", True)
            max_resultados = st.session_state.get("max_resultados_fp", 50)

            # ── 1. Lookup no Excel ────────────────────────────────────────────
            fontes_rag = []
            resultado_matriz = ""
            if usar_fp_matriz:
                csv_path   = os.path.join(BASE_DOCS_DIR, MATRIX_CSV)
                excel_path = os.path.join(BASE_DOCS_DIR, MATRIX_EXCEL)
                if not os.path.exists(csv_path) and not os.path.exists(excel_path):
                    st.error(
                        f"Nenhum arquivo de matriz encontrado em base_docs/. "
                        f"Coloque {MATRIX_CSV} (preferido) ou {MATRIX_EXCEL} lá e reinicie o app."
                    )
                    st.stop()
                df_fp = carregar_df_fp_matriz(excel_path, SHEET_FP_MATRIZ)

                # Pergunta de contagem — responde diretamente com len(df)
                if _detectar_pergunta_contagem(input_usuario):
                    resultado_matriz = f"O portfólio FleetPro possui **{len(df_fp)} itens** cadastrados na Matriz FP."

                tt_code = detectar_tt_code(input_usuario)
                colunas_equip = detectar_busca_por_equipamento(input_usuario)
                termo_marketing = detectar_busca_marketing(input_usuario)

                if _detectar_pergunta_contagem(input_usuario):
                    pass  # já tratado acima
                elif tt_code:
                    # Busca por TT code
                    resultado_matriz = buscar_por_tt_code(df_fp, tt_code, max_resultados)
                elif colunas_equip and termo_marketing:
                    # ── Busca combinada: equipamento + categoria ──────────────
                    # Ex: "rolamento para trator case" → filtra CASE IH + ROLAMENTOS
                    resultado_matriz = buscar_equip_e_marketing(
                        df_fp, input_usuario, colunas_equip, termo_marketing, max_resultados
                    )
                elif colunas_equip:
                    # Busca por fabricante / tipo de máquina
                    resultado_matriz = buscar_por_equipamento(df_fp, input_usuario, colunas_equip, max_resultados)
                elif termo_marketing:
                    # Busca por categoria de produto na coluna MARKETING PROJECT
                    resultado_matriz = buscar_por_marketing(df_fp, termo_marketing, input_usuario, max_resultados)
                else:
                    resultado_matriz = procurar_pn(df_fp, input_usuario, max_resultados)

                # ── Fallback: busca por texto livre na coluna DESCRIPTION ──────
                # Ativado quando nenhuma busca retornou resultado útil
                if not resultado_matriz or resultado_matriz.startswith("Não encontrei") or resultado_matriz.startswith("PN não encontrado"):
                    resultado_descricao = buscar_por_description(df_fp, input_usuario, max_resultados)
                    if resultado_descricao:
                        resultado_matriz = resultado_descricao

            # ── 2. Busca RAG (documentos locais + site FleetPro) ──────────────
            contexto_rag = ""
            fontes_rag = []
            if usar_rag and chat_model is not None and not st.session_state.get("_rag_indisponivel"):
                try:
                    vs = obter_vectorstore()
                    contexto_rag, fontes_rag = buscar_no_rag(vs, input_usuario)
                except Exception:
                    contexto_rag = ""
                    fontes_rag = []

            # ── 3. Montar resposta ────────────────────────────────────────────
            if usar_fp_matriz and not usar_rag and chat_model is None:
                st.markdown(resultado_matriz)
                resposta = resultado_matriz

            elif chat_model is not None:
                # Bloqueia LLM se não há nenhum dado — evita alucinação
                if not resultado_matriz and not contexto_rag:
                    resposta = "Não encontrei informações sobre isso na base FleetPro. Tente buscar por PN, categoria de produto ou equipamento."
                    st.markdown(resposta)
                    memoria.chat_memory.add_user_message(input_usuario)
                    memoria.chat_memory.add_ai_message(resposta)
                    st.session_state["memoria"] = memoria
                    st.stop()

                blocos = []

                if resultado_matriz:
                    # Trunca para não explodir o contexto
                    resultado_matriz_trunc = resultado_matriz[:2000] + ("\n...(truncado)" if len(resultado_matriz) > 2000 else "")
                    blocos.append(
                        "## Resultado da Busca na Matriz FP (dados do Excel)\n\n"
                        + resultado_matriz_trunc
                    )
                    # Adiciona matriz como fonte visível
                    fontes_rag.insert(0, {
                        "label": "📊 Matriz FP — resultado encontrado",
                        "path": "",
                        "tipo": "matriz"
                    })
                    # Adiciona Matriz FP como fonte consultada
                    fontes_rag.append({
                        "label": "📊 Matriz FP (Excel) — " + MATRIX_EXCEL,
                        "path": os.path.join(BASE_DOCS_DIR, MATRIX_EXCEL),
                        "tipo": "matriz_excel"
                    })

                if contexto_rag:
                    blocos.append(
                        "## Conhecimento Adicional (guia de objeções / recomendações / site FleetPro)\n\n"
                        + contexto_rag
                    )

                if blocos:
                    contexto_completo = "\n\n---\n\n".join(blocos)
                    perfil = st.session_state.get("perfil_usuario")

                    if perfil is None:
                        st.session_state["pergunta_pendente"] = input_usuario
                        st.session_state["contexto_pendente"] = contexto_completo
                        resposta = (
                            "Antes de responder, preciso entender melhor como posso te ajudar! 😊\n\n"
                            "**Você é:**\n"
                            "- **1️⃣ Vendedor** — quero argumentos para atender meu cliente no balcão\n"
                            "- **2️⃣ Usuário da peça** — quero saber se o FleetPro é a melhor opção para minha máquina\n\n"
                            "_Digite o número ou o nome da opção._"
                        )
                        st.markdown(resposta)
                        memoria.chat_memory.add_user_message(input_usuario)
                        memoria.chat_memory.add_ai_message(resposta)
                        st.session_state["memoria"] = memoria
                        st.stop()

                    elif perfil == "vendedor":
                        # ── Busca simples: mostra peças + botão de argumentação ──
                        if _eh_busca_simples(input_usuario) and resultado_matriz and not contexto_rag:
                            resposta = resultado_matriz
                            st.markdown(resposta)
                            if st.button(
                                "💬 Ver argumentação de venda",
                                key=f"btn_arg_{hash(input_usuario)}",
                                type="secondary",
                            ):
                                st.session_state["_argumentacao_pendente"] = {
                                    "input_usuario": input_usuario,
                                    "resultado_matriz": resultado_matriz,
                                    "contexto_rag": contexto_rag,
                                }
                                st.rerun()
                            memoria.chat_memory.add_user_message(input_usuario)
                            memoria.chat_memory.add_ai_message(resposta)
                            st.session_state["memoria"] = memoria
                            st.stop()

                        prompt = (
                            f"Você é o FleetPro Expert, um assistente de vendas especializado em apoiar vendedores de balcão "
                            f"na comercialização de peças de reposição FleetPro para máquinas agrícolas.\n\n"
                            f"Seu papel é ajudar o VENDEDOR a:\n"
                            f"- Identificar rapidamente o PN FleetPro correto para o cliente\n"
                            f"- Ter argumentos técnicos e comerciais prontos para contornar objeções\n"
                            f"- Sempre oferecer a possibilidade de outros itens em conjunto com a demanda (Cross Selling)\n"
                            f"- Explicar as vantagens do FleetPro frente à peça original ou concorrente\n"
                            f"- Transmitir confiança e credibilidade na indicação do produto\n\n"
                            f"Tom das respostas:\n"
                            f"- Direto e objetivo — o vendedor está atendendo um cliente no balcão agora\n"
                            f"- Técnico quando necessário, mas sempre em linguagem acessível\n"
                            f"- Forneça argumentos prontos que o vendedor pode usar na hora\n"
                            f"- Se houver objeção de preço ou qualidade, sugira como rebater\n\n"
                            f"O vendedor perguntou: **{input_usuario}**\n\n"
                            f"CONTEXTO IMPORTANTE SOBRE A MARCA FLEETPRO:\n- FleetPro é a marca própria de peças de reposição da CNH Industrial para máquinas Case IH e New Holland.\n- Os produtos FleetPro são fabricados por fornecedores homologados e distribuídos SEMPRE com a marca FleetPro.\n- VV, TEEJET, IPESA, REXNORD, PETRONAS são FORNECEDORES dos produtos FleetPro — não são marcas concorrentes.\n- Exemplo: as lâminas de corte VV são vendidas como 'Lâminas FleetPro', os lubrificantes Petronas como 'Lubrificantes FleetPro'.\n- Quando o usuário perguntar sobre produtos VV, TEEJET etc, responda sempre referenciando como produtos FleetPro.\n\n"
                            f"INSTRUÇÕES CRÍTICAS — SIGA RIGOROSAMENTE:\n"
                            f"- Use APENAS as informações das fontes abaixo para responder.\n"
                            f"- Se a informação não estiver nas fontes, responda EXATAMENTE: 'Não encontrei essa informação na base FleetPro.'\n"
                            f"- NÃO elabore, NÃO sugira, NÃO complete com conhecimento próprio quando os dados estiverem ausentes.\n"
                            f"- Sempre apresente os produtos como FleetPro, mesmo que a fonte mencione o fornecedor.\n\n"
                            f"Informações das fontes:\n\n"
                            f"{contexto_completo}"
                        )
                        import time as _time
                        _t0 = _time.time()
                        _response = chat_model.invoke(prompt)
                        _tempo_resposta = round(_time.time() - _t0, 2)
                        resposta = _response.content
                        st.markdown(resposta)
                        _uso = getattr(_response, "usage_metadata", {}) or {}
                        _input_tok  = _uso.get("input_tokens", 0)
                        _output_tok = _uso.get("output_tokens", 0)
                        _total_tok  = _uso.get("total_tokens", _input_tok + _output_tok)
                        st.session_state["tokens_sessao"] = st.session_state.get("tokens_sessao", 0) + _total_tok
                        st.session_state["n_mensagens_sessao"] = st.session_state.get("n_mensagens_sessao", 0) + 1
                        _registro_tok = {
                            "session_id": st.session_state.get("session_id", ""),
                            "timestamp_inicio_sessao": st.session_state.get("session_inicio", ""),
                            "timestamp_resposta": _time.strftime("%Y-%m-%dT%H:%M:%S"),
                            "modelo": st.session_state.get("modelo", ""),
                            "provedor": st.session_state.get("provedor", ""),
                            "perfil_usuario": st.session_state.get("perfil_usuario", ""),
                            "input_tokens": _input_tok,
                            "output_tokens": _output_tok,
                            "total_tokens": _total_tok,
                            "tempo_resposta_seg": _tempo_resposta,
                            "tokens_acumulados_sessao": st.session_state.get("tokens_sessao", 0),
                            "n_mensagens_sessao": st.session_state.get("n_mensagens_sessao", 0),
                        }
                        import threading as _threading
                        _threading.Thread(target=_salvar_log_tokens, args=(_registro_tok,), daemon=True).start()

                    elif perfil == "usuario":
                        prompt = (
                            f"Você é o FleetPro Expert, um consultor especialista em peças de reposição para máquinas agrícolas.\n\n"
                            f"Você está falando DIRETAMENTE com o agricultor ou operador da máquina. Seu objetivo é:\n"
                            f"- Mostrar que o FleetPro é a melhor escolha para a máquina dele\n"
                            f"- Destacar que a qualidade é equivalente ou superior à peça original\n"
                            f"- Ressaltar economia de custo sem abrir mão de desempenho\n"
                            f"- Transmitir segurança e confiança na marca FleetPro\n"
                            f"- Incentivar o cliente a pedir o FleetPro ao seu revendedor\n\n"
                            f"Tom das respostas:\n"
                            f"- Próximo, empático e direto — fale como um especialista de confiança\n"
                            f"- Valorize a economia e a produtividade que o FleetPro proporciona\n"
                            f"- Use exemplos práticos do dia a dia do campo quando possível\n"
                            f"- Finalize sempre incentivando a buscar o FleetPro no revendedor mais próximo\n\n"
                            f"O usuário perguntou: **{input_usuario}**\n\n"
                            f"CONTEXTO IMPORTANTE SOBRE A MARCA FLEETPRO:\n- FleetPro é a marca própria de peças de reposição da CNH Industrial para máquinas Case IH e New Holland.\n- Os produtos FleetPro são fabricados por fornecedores homologados e distribuídos SEMPRE com a marca FleetPro.\n- VV, TEEJET, IPESA, REXNORD, PETRONAS são FORNECEDORES dos produtos FleetPro — não são marcas concorrentes.\n- Exemplo: as lâminas de corte VV são vendidas como 'Lâminas FleetPro', os lubrificantes Petronas como 'Lubrificantes FleetPro'.\n- Quando o usuário perguntar sobre produtos VV, TEEJET etc, responda sempre referenciando como produtos FleetPro.\n\n"
                            f"INSTRUÇÕES CRÍTICAS — SIGA RIGOROSAMENTE:\n"
                            f"- Use APENAS as informações das fontes abaixo para responder.\n"
                            f"- Se a informação não estiver nas fontes, responda EXATAMENTE: 'Não encontrei essa informação na base FleetPro.'\n"
                            f"- NÃO elabore, NÃO sugira, NÃO complete com conhecimento próprio quando os dados estiverem ausentes.\n"
                            f"- Sempre apresente os produtos como FleetPro, mesmo que a fonte mencione o fornecedor.\n\n"
                            f"Informações das fontes:\n\n"
                            f"{contexto_completo}"
                        )
                        import time as _time
                        _t0 = _time.time()
                        _response = chat_model.invoke(prompt)
                        _tempo_resposta = round(_time.time() - _t0, 2)
                        resposta = _response.content
                        st.markdown(resposta)
                        _uso = getattr(_response, "usage_metadata", {}) or {}
                        _input_tok  = _uso.get("input_tokens", 0)
                        _output_tok = _uso.get("output_tokens", 0)
                        _total_tok  = _uso.get("total_tokens", _input_tok + _output_tok)
                        st.session_state["tokens_sessao"] = st.session_state.get("tokens_sessao", 0) + _total_tok
                        st.session_state["n_mensagens_sessao"] = st.session_state.get("n_mensagens_sessao", 0) + 1
                        _registro_tok = {
                            "session_id": st.session_state.get("session_id", ""),
                            "timestamp_inicio_sessao": st.session_state.get("session_inicio", ""),
                            "timestamp_resposta": _time.strftime("%Y-%m-%dT%H:%M:%S"),
                            "modelo": st.session_state.get("modelo", ""),
                            "provedor": st.session_state.get("provedor", ""),
                            "perfil_usuario": st.session_state.get("perfil_usuario", ""),
                            "input_tokens": _input_tok,
                            "output_tokens": _output_tok,
                            "total_tokens": _total_tok,
                            "tempo_resposta_seg": _tempo_resposta,
                            "tokens_acumulados_sessao": st.session_state.get("tokens_sessao", 0),
                            "n_mensagens_sessao": st.session_state.get("n_mensagens_sessao", 0),
                        }
                        import threading as _threading
                        _threading.Thread(target=_salvar_log_tokens, args=(_registro_tok,), daemon=True).start()

                else:
                    perfil = st.session_state.get("perfil_usuario")
                    perfil_texto = ""
                    if perfil == "vendedor":
                        perfil_texto = "Você está auxiliando um VENDEDOR DE BALCÃO."
                    elif perfil == "usuario":
                        perfil_texto = "Você está falando DIRETAMENTE com um agricultor ou operador de máquina."
                    prompt = (
                        f"Você é o FleetPro Expert, assistente especializado em peças de reposição FleetPro "
                        f"para máquinas agrícolas das marcas Case IH, New Holland, John Deere e outras.\n"
                        f"A FleetPro oferece peças de qualidade equivalente à original com custo mais competitivo.\n"
                        f"CONTEXTO IMPORTANTE: VV, TEEJET, IPESA, REXNORD e PETRONAS são FORNECEDORES dos produtos FleetPro.\n"
                        f"Esses fornecedores fabricam os produtos que são distribuídos SEMPRE com a marca FleetPro pela rede CNH.\n"
                        f"Nunca trate esses nomes como marcas concorrentes — são parceiros que compõem o portfólio FleetPro.\n\n"
                        f"{perfil_texto}\n\n"
                        f"Responda SEMPRE como FleetPro Expert. Nunca diga que é um modelo de IA genérico.\n"
                        f"Se não souber algo específico sobre um PN, oriente o usuário a buscar na matriz ou "
                        f"perguntar ao revendedor FleetPro mais próximo.\n\n"
                        f"Pergunta: {input_usuario}"
                    )
                    import time as _time
                    _t0 = _time.time()
                    _response = chat_model.invoke(prompt)
                    _tempo_resposta = round(_time.time() - _t0, 2)
                    resposta = _response.content
                    st.markdown(resposta)
                    _uso = getattr(_response, "usage_metadata", {}) or {}
                    _input_tok  = _uso.get("input_tokens", 0)
                    _output_tok = _uso.get("output_tokens", 0)
                    _total_tok  = _uso.get("total_tokens", _input_tok + _output_tok)
                    st.session_state["tokens_sessao"] = st.session_state.get("tokens_sessao", 0) + _total_tok
                    st.session_state["n_mensagens_sessao"] = st.session_state.get("n_mensagens_sessao", 0) + 1
                    _registro_tok = {
                        "session_id": st.session_state.get("session_id", ""),
                        "timestamp_inicio_sessao": st.session_state.get("session_inicio", ""),
                        "timestamp_resposta": _time.strftime("%Y-%m-%dT%H:%M:%S"),
                        "modelo": st.session_state.get("modelo", ""),
                        "provedor": st.session_state.get("provedor", ""),
                        "perfil_usuario": st.session_state.get("perfil_usuario", ""),
                        "input_tokens": _input_tok,
                        "output_tokens": _output_tok,
                        "total_tokens": _total_tok,
                        "tempo_resposta_seg": _tempo_resposta,
                        "tokens_acumulados_sessao": st.session_state.get("tokens_sessao", 0),
                        "n_mensagens_sessao": st.session_state.get("n_mensagens_sessao", 0),
                    }
                    import threading as _threading
                    _threading.Thread(target=_salvar_log_tokens, args=(_registro_tok,), daemon=True).start()

            else:
                if resultado_matriz:
                    st.markdown(resultado_matriz)
                    resposta = resultado_matriz
                else:
                    resposta = (
                        "⚙️ Nenhum modelo de linguagem inicializado.\n\n"
                        "Para perguntas e recomendações, configure o LLM na sidebar "
                        "(**Modelo de Linguagem**) e clique em **Inicializar Agente**.\n\n"
                        "Para buscar um PN diretamente, basta digitar o número da peça."
                    )
                    st.markdown(resposta)

        except Exception as e:
            st.error("Erro ao responder.")
            st.exception(e)
            st.stop()

    memoria.chat_memory.add_user_message(input_usuario)
    memoria.chat_memory.add_ai_message(resposta)
    st.session_state["memoria"] = memoria

    # ── Exibe fontes consultadas ──────────────────────────────────────────
    if fontes_rag:
        with st.expander("\U0001F4DA Fontes consultadas (" + str(len(fontes_rag)) + ")", expanded=False):
            for fonte in fontes_rag:
                tipo = fonte.get("tipo", "")
                label = fonte.get("label", "")
                path = fonte.get("path", "")
                st.markdown("- " + label)


# ======================
# UI – Barra lateral
# ======================
def sidebar():
    st.title("⚙️ Configurações")
    st.image("base_docs/cnh_logo.png", width=180)

    # ── Seção 1: Modelo de linguagem (LLM) ──────────────────────────────────
    with st.expander("🤖 Modelo de Linguagem (LLM)", expanded=True):
        st.caption(
            "O LLM combina os resultados da busca no Excel com o conhecimento dos documentos "
            "e gera uma resposta enriquecida. Se não inicializado, o app retorna apenas os dados brutos da matriz."
        )

        provedor = st.selectbox("Provedor", list(CONFIG_MODELOS.keys()), key="sel_provedor")
        modelo = st.selectbox("Modelo", CONFIG_MODELOS[provedor]["modelos"], key="sel_modelo")

        # Lê cookie salvo para pré-preencher o campo
        _cookie_key = f"fp_apikey_{provedor.lower()}"
        _saved_key = ""
        if _COOKIES_OK:
            try:
                _saved_key = _cookie_manager.get(_cookie_key) or ""
            except Exception:
                _saved_key = ""

        api_key = st.text_input(
            f"API key ({provedor})",
            value=st.session_state.get(f"api_key_{provedor}", _saved_key),
            type="password",
            key="input_api_key_llm",
        )
        st.session_state[f"api_key_{provedor}"] = api_key

        if provedor == "OpenAI" and api_key:
            st.session_state["api_key_openai_rag"] = api_key

        if st.button("🚀 Inicializar FleetPro_Expert", use_container_width=True):
            inicializar_FleetPro(provedor, modelo, api_key)
            # Salva a key no cookie (expira em 30 dias)
            if _COOKIES_OK and api_key:
                try:
                    _cookie_manager.set(_cookie_key, api_key, max_age=60*60*24*30)
                except Exception:
                    pass

        st.divider()
        chat = st.session_state.get("chat")
        if chat:
            st.success(
                f"✅ Ativo: {st.session_state.get('provedor')} / {st.session_state.get('modelo')}"
            )
        else:
            st.warning("⚠️ Nenhum modelo inicializado. O chat usará apenas o lookup direto.")

    # ── Seção 2: FP MATRIZ ───────────────────────────────────────────────────
    with st.expander("📊 Lookup FP MATRIZ (Excel)", expanded=False):
        usar_fp_matriz = st.toggle(
            "Buscar PN na Matriz FP",
            value=st.session_state.get("usar_fp_matriz", True),
        )
        st.session_state["usar_fp_matriz"] = usar_fp_matriz

        max_resultados_fp = st.slider(
            "Máximo de resultados retornados",
            min_value=5, max_value=200,
            value=st.session_state.get("max_resultados_fp", 50),
        )
        st.session_state["max_resultados_fp"] = max_resultados_fp

        st.caption(
            "A busca normaliza o PN (remove espaços, hífens e pontos; maiúsculo) "
            "e compara com PN GEN, PN ALTERNATIVE, PN NXP e PN FLEETPRO."
        )

    # ── Seção 3: RAG / Documentos de conhecimento ────────────────────────────
    with st.expander("📂 Documentos de Conhecimento (RAG)", expanded=False):
        st.info(
            "Coloque seus documentos (.pdf, .txt, .csv, .pptx, .md, .docx) na pasta `base_docs/` ou em subpastas. O sistema indexa automaticamente todos os arquivos e subpastas. "
            "O sistema irá indexá-los e usará o conteúdo para ajudar com objeções, "
            "recomendações e dúvidas técnicas."
        )

        usar_rag = st.toggle(
            "Usar documentos de conhecimento no chat",
            value=st.session_state.get("usar_rag", True),
        )
        st.session_state["usar_rag"] = usar_rag

        if st.session_state.get("_rag_indisponivel"):
            st.warning("⚠️ RAG indisponível — modelo de embeddings não pôde ser carregado (sem acesso à internet ou rede bloqueada). O app está rodando **sem RAG**, usando apenas a Matriz FP.")
        else:
            st.caption("✅ Embeddings gratuitos (HuggingFace) — sem necessidade de API key.")

        col1, col2 = st.columns(2)
        with col1:
            if st.button("🔄 Reindexar", use_container_width=True):
                try:
                    resetar_vectorstore()
                    vs = obter_vectorstore()
                    if vs:
                        st.success("Indexado com sucesso!")
                    else:
                        st.warning("Nenhum documento encontrado.")
                except Exception as e:
                    st.error(f"Erro: {e}")

        with col2:
            if st.button("📋 Ver docs", use_container_width=True):
                arquivos = _listar_arquivos_rag(BASE_DOCS_DIR)
                if arquivos:
                    st.success(f"{len(arquivos)} arquivo(s) encontrado(s):")
                    for a in arquivos:
                        rel = os.path.relpath(a, BASE_DOCS_DIR)
                        st.caption(f"• {rel}")
                else:
                    st.warning("Nenhum documento local encontrado.")




# ======================
# Salvar erro no GitHub
# ======================
def _salvar_erro_github(registro: dict):
    """Commita erros_reportados.csv no repositório GitHub via API."""
    import base64, json, csv, io, urllib.request, urllib.error

    token = st.secrets.get("GITHUB_TOKEN", os.environ.get("GITHUB_TOKEN", ""))
    if not token:
        raise RuntimeError("GITHUB_TOKEN não configurado nos secrets do Streamlit.")

    owner = "Antoniopneto9"
    repo  = "Agente_FleetPro-1"
    path  = "erros_reportados.csv"
    api   = f"https://api.github.com/repos/{owner}/{repo}/contents/{path}"
    gh_headers = {
        "Authorization": f"token {token}",
        "Accept": "application/vnd.github+json",
        "Content-Type": "application/json",
    }

    def _fetch():
        req = urllib.request.Request(api + "?ref=prod", headers=gh_headers)
        try:
            with urllib.request.urlopen(req) as resp:
                data = json.loads(resp.read())
                sha = data["sha"]
                # content pode ter quebras de linha — remover antes de decodificar
                content_b64 = data["content"].replace("\n", "")
                return sha, base64.b64decode(content_b64).decode("utf-8")
        except urllib.error.HTTPError as e:
            if e.code == 404:
                return None, ""
            raise

    def _montar_conteudo(conteudo_atual):
        campos = list(registro.keys())
        buf = io.StringIO()
        writer = csv.DictWriter(buf, fieldnames=campos, quoting=csv.QUOTE_ALL)
        writer.writeheader()
        if conteudo_atual.strip():
            reader = csv.DictReader(io.StringIO(conteudo_atual))
            for row in reader:
                linha = {k: row.get(k, "") for k in campos}
                writer.writerow(linha)
        writer.writerow(registro)
        return buf.getvalue()

    def _put(sha, novo_conteudo):
        payload = {
            "message": f"erro reportado: {registro['timestamp'][:10]}",
            "content": base64.b64encode(novo_conteudo.encode("utf-8")).decode("utf-8"),
            "branch": "prod",
        }
        if sha:
            payload["sha"] = sha
        req = urllib.request.Request(api, data=json.dumps(payload).encode(), headers=gh_headers, method="PUT")
        try:
            with urllib.request.urlopen(req) as resp:
                resp.read()
        except urllib.error.HTTPError as e:
            raise RuntimeError(f"GitHub API error {e.code}: {e.read().decode('utf-8', errors='replace')}")

    sha, conteudo_atual = _fetch()
    _put(sha, _montar_conteudo(conteudo_atual))


# ======================
# Log de Tokens
# ======================
def _salvar_log_tokens(registro: dict):
    """Commita uso_tokens.csv no repositório GitHub via API (background thread)."""
    import base64, json, csv, io, urllib.request, urllib.error

    token = st.secrets.get("GITHUB_TOKEN", os.environ.get("GITHUB_TOKEN", ""))
    if not token:
        return  # sem token, silencioso — não bloquear a UI

    owner = "Antoniopneto9"
    repo  = "Agente_FleetPro-1"
    path  = "uso_tokens.csv"
    api   = f"https://api.github.com/repos/{owner}/{repo}/contents/{path}"
    gh_headers = {
        "Authorization": f"token {token}",
        "Accept": "application/vnd.github+json",
        "Content-Type": "application/json",
    }

    def _fetch():
        req = urllib.request.Request(api + "?ref=prod", headers=gh_headers)
        try:
            with urllib.request.urlopen(req) as resp:
                data = json.loads(resp.read())
                sha = data["sha"]
                content_b64 = data["content"].replace("\n", "")
                return sha, base64.b64decode(content_b64).decode("utf-8")
        except urllib.error.HTTPError as e:
            if e.code == 404:
                return None, ""
            raise

    def _montar_conteudo(conteudo_atual):
        campos = list(registro.keys())
        buf = io.StringIO()
        writer = csv.DictWriter(buf, fieldnames=campos, quoting=csv.QUOTE_ALL)
        writer.writeheader()
        if conteudo_atual.strip():
            reader = csv.DictReader(io.StringIO(conteudo_atual))
            for row in reader:
                linha = {k: row.get(k, "") for k in campos}
                writer.writerow(linha)
        writer.writerow(registro)
        return buf.getvalue()

    def _put(sha, novo_conteudo):
        payload = {
            "message": f"log tokens: {registro.get('timestamp_resposta', '')[:10]}",
            "content": base64.b64encode(novo_conteudo.encode("utf-8")).decode("utf-8"),
            "branch": "prod",
        }
        if sha:
            payload["sha"] = sha
        req = urllib.request.Request(api, data=json.dumps(payload).encode(), headers=gh_headers, method="PUT")
        try:
            with urllib.request.urlopen(req) as resp:
                resp.read()
        except urllib.error.HTTPError:
            pass  # silencioso em background

    try:
        sha, conteudo_atual = _fetch()
        _put(sha, _montar_conteudo(conteudo_atual))
    except Exception:
        pass  # nunca deixar o background thread crashar a UI


# ======================
# UI – Reportar Erro
# ======================
def popup_feedback():
    """Captura descrição do erro + histórico do chat principal e salva no projeto."""
    import datetime, csv

    if "_feedback_reset_count" not in st.session_state:
        st.session_state["_feedback_reset_count"] = 0

    with st.sidebar:
        st.divider()
        with st.expander("🐛 Reportar Erro", expanded=False):
            reset_count = st.session_state["_feedback_reset_count"]
            descricao = st.text_area(
                "O que aconteceu?",
                placeholder="Ex: Perguntei sobre o PN X e o modelo retornou Y errado.",
                height=100,
                key=f"feedback_descricao_{reset_count}",
            )

            if st.button("📤 Enviar Erro", use_container_width=True, key=f"btn_enviar_feedback_{reset_count}", type="primary"):
                if not descricao.strip():
                    st.warning("Descreva o erro antes de enviar.")
                else:
                    memoria: ConversationBufferMemory = st.session_state.get("memoria", ConversationBufferMemory())
                    historico_str = " | ".join(
                        f"{m.type.upper()}: {m.content.replace(chr(10), ' ').replace(chr(13), ' ')}"
                        for m in memoria.buffer_as_messages
                    )

                    registro = {
                        "timestamp": datetime.datetime.now().isoformat(),
                        "descricao": descricao.strip(),
                        "modelo": st.session_state.get("modelo", "N/A"),
                        "provedor": st.session_state.get("provedor", "N/A"),
                        "historico_chat": historico_str,
                    }

                    try:
                        _salvar_erro_github(registro)
                        st.session_state["memoria"] = ConversationBufferMemory()
                        st.session_state["mensagens"] = []
                        st.session_state["_feedback_reset_count"] += 1
                        st.rerun()
                    except Exception as e:
                        st.error(f"Erro ao salvar: {e}")


# ======================
# Main
# ======================
def main():
    pagina_chat()
    with st.sidebar:
        sidebar()
        popup_feedback()


if __name__ == "__main__":
    safe_run(main)